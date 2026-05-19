"""Render a SlideDeck to .pptx — academic defense style (v8).

Design: monochrome cream bg (#FBFAF7), near-black text, no accent colors,
formal graduation-defense aesthetic, compact layout, no brand tag.
Template-swap: caller may supply a .pptx master via `template_path`.

v8 changes vs v7:
- section_divider redesigned: left vertical title block + right key-points card
  (Design B: 左定锚 + 右卡片, meeting/report feel)

v7 changes vs v6:
- Outline slide now shows 4-5 grouped sections (outline_grouped kind)
- section_divider replaces per-chapter dividers; absorbs pure-bullet chapter content
- closing_rich slide uses LLM paper_brief (5-7 bullets + takeaway)
- _dispatch handles "section_divider", "outline_grouped", "closing_rich"
"""
from __future__ import annotations

import datetime
import re
from pathlib import Path
from typing import ClassVar

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

from stages.s09_render.model import Document
from stages.s09_render.renderers import RENDERERS
from stages.s09_render.renderers.base import Renderer
from stages.s09_render.slide_planner import Slide, SlideDeck, SlidePlanner


class _T:
    """Design token namespace — monochrome academic palette."""
    W, H         = Inches(13.333), Inches(7.5)
    BG           = RGBColor(0xFB, 0xFA, 0xF7)   # very light cream
    TEXT         = RGBColor(0x1A, 0x1A, 0x1A)   # near-black body
    TEXT_DIM     = RGBColor(0x66, 0x66, 0x66)   # subtitle / metadata
    TEXT_FAINT   = RGBColor(0xA8, 0xA8, 0xA8)   # footer / fine rules
    RULE         = RGBColor(0xD0, 0xD0, 0xD0)   # thin separator lines
    EA_SERIF     = "Source Han Serif SC"
    EA_SANS      = "PingFang SC"
    LAT_SERIF    = "Cambria"
    LAT_SANS     = "Calibri"


_IDX_TITLE = 0
_IDX_BLANK = 6

_ZH_NUMS = ["一","二","三","四","五","六","七","八","九","十",
             "十一","十二","十三","十四","十五","十六","十七","十八","十九","二十"]

_MARKERS = ["❶", "❷", "❸", "❹", "❺", "❻", "❼", "❽", "❾", "❿"]

_GRAY88 = RGBColor(0x88, 0x88, 0x88)


class _S:
    """Localized UI string pairs — (zh, en).

    v13: zh values are PURE Chinese; en values are PURE English. No mixed strings.
    """
    OUTLINE_TITLE   = ("目录",          "Contents")
    SECTION_EYEBROW = ("本节要点",       "Key Points")
    NAV_HINT        = ("» 继续 →",      "» continue →")
    FIG_EYEBROW     = ("深度观察",       "Key Observations")
    COMBINED_KW     = ("要点",          "Key Points")
    CLOSING_TITLE   = ("结论",          "Conclusion")
    CHAPTER_FOOTER  = ("章节：",        "Chapter: ")
    BYLINE_LABEL    = ("演讲者  ·  单位  ·  日期", "PRESENTER  ·  AFFILIATION  ·  DATE")

    @staticmethod
    def pick(pair: tuple[str, str], lang: str) -> str:
        return pair[0] if lang == "zh" else pair[1]


class PptxRenderer(Renderer):
    extension: ClassVar[str] = "pptx"

    def __init__(self, summaries: dict | None = None,
                 outline: list[dict] | None = None,
                 paper_brief: dict | None = None,
                 template_path: Path | None = None,
                 presenter: str | None = None,
                 affiliation: str | None = None,
                 subtitle: str | None = None):
        self.summaries = summaries
        self.outline = outline
        self.paper_brief = paper_brief
        self.template_path = template_path
        self.presenter = presenter or "Paper2MD Auto-Analysis"
        self.affiliation = affiliation or "Scientific Paper Deep Analysis"
        self.subtitle = subtitle
        self._ch: int = -1          # chapter index, incremented on each divider
        self._sec_idx: int = -1     # section_divider index (v7)
        self._fig_idx: int = 0      # figure counter for alternating layout
        self._cur_chapter: str = "" # current chapter heading for footer
        self._cur_section: str = "" # current section name for footer (v7)

    def render(self, doc: Document, out_path: Path) -> None:
        deck = SlidePlanner(lang=doc.lang).plan(
            doc, self.summaries,
            outline=self.outline,
            paper_brief=self.paper_brief,
        )
        if self.template_path is not None:
            prs = Presentation(str(self.template_path))
        else:
            prs = Presentation()
            prs.slide_width, prs.slide_height = _T.W, _T.H
        total = len(deck.slides)
        self._ch = -1
        self._sec_idx = -1
        self._fig_idx = 0
        self._cur_chapter = ""
        self._cur_section = ""
        # Build heading → 0-based chapter index lookup for §N tracking
        self._ch_idx: dict[str, int] = {
            ch.heading: i for i, ch in enumerate(doc.chapters)
        }
        for idx, slide in enumerate(deck.slides, 1):
            self._dispatch(prs, slide, idx, total, doc)
        prs.save(str(out_path))

    def _dispatch(self, prs, slide, idx, total, doc):
        kw = dict(idx=idx, total=total, doc=doc)
        if   slide.kind == "title":          self._title(prs, slide, **kw)
        elif slide.kind == "outline":        self._outline(prs, slide, **kw)
        elif slide.kind == "outline_grouped": self._outline_grouped(prs, slide, **kw)
        elif slide.kind == "section_divider":
            self._sec_idx += 1
            self._cur_section = slide.title
            self._section_divider(prs, slide, **kw)
        elif slide.kind == "divider":
            self._ch += 1
            self._cur_chapter = slide.title
            self._divider(prs, slide, **kw)
        elif slide.kind == "figure":
            self._figure(prs, slide, **kw)
            self._fig_idx += 1
        elif slide.kind == "combined":
            self._combined(prs, slide, **kw)
            self._fig_idx += 1
        elif slide.kind == "closing_rich":
            self._closing_rich(prs, slide, **kw)
        else:
            if slide.kind in ("bullets", "closing"):
                self._cur_chapter = slide.title
            self._bullets(prs, slide, **kw)

    # ── slide builders ─────────────────────────────────────────────────────────

    def _title(self, prs, slide, *, idx, total, doc):
        """Title slide — v13: tighter vertical centering, pure-language byline."""
        s = prs.slides.add_slide(self._lay(prs, _IDX_TITLE))
        _bg(s)
        T = _T
        title_len = len(slide.title)
        title_pt = 30 if title_len < 60 else (26 if title_len < 100 else 22)

        ph = s.shapes.title
        if ph is not None:
            ph.left   = Inches(1.2)
            ph.top    = Inches(1.5)      # v13: was 1.0, move up by raising relative position
            ph.width  = Inches(10.9)
            ph.height = Inches(2.5)      # v13: was 2.8, tighter
            ph.text_frame.word_wrap = True
            ph.text_frame.text = slide.title
            p = ph.text_frame.paragraphs[0]
            if p.runs:
                _run_style(p.runs[0], Pt(title_pt), True, T.TEXT, T.LAT_SERIF, T.EA_SERIF)
            p.alignment = PP_ALIGN.CENTER

        for ph2 in list(s.placeholders):
            if ph2.placeholder_format.idx == 1:
                ph2._element.getparent().remove(ph2._element); break

        rule_w = Inches(3.5)
        rule_x = (_T.W - rule_w) / 2
        _line(s, rule_x, Inches(4.1), rule_x + rule_w, Inches(4.1), T.RULE, Pt(0.75))

        subtitle_text = self.subtitle or ""
        if subtitle_text:
            _tb1(s, subtitle_text,
                 Inches(1.5), Inches(4.25), Inches(10.333), Inches(0.45),
                 Pt(13), T.TEXT_DIM, T.LAT_SANS, T.EA_SANS, italic=True, align=PP_ALIGN.CENTER)
            meta_top = Inches(4.75)
        else:
            meta_top = Inches(4.3)

        today = datetime.date.today()
        date_str = f"{today.year}-{today.month:02d}-{today.day:02d}"
        byline_x = Inches(1.5)
        byline_w = Inches(10.333)

        # v13: pure-language byline label
        byline_label = _S.pick(_S.BYLINE_LABEL, doc.lang)
        _tb1(s, byline_label,
             byline_x, meta_top, byline_w, Inches(0.28),
             Pt(8), T.TEXT_FAINT, T.LAT_SANS, T.EA_SANS, align=PP_ALIGN.CENTER)

        _tb1(s, f"{self.presenter}  ·  {self.affiliation}  ·  {date_str}",
             byline_x, meta_top + Inches(0.28), byline_w, Inches(0.42),
             Pt(14), T.TEXT_DIM, T.LAT_SANS, T.EA_SANS, bold=False, align=PP_ALIGN.CENTER)

        _footer(s, idx, total, chapter="", lang=doc.lang)
        _notes(s, slide.notes)

    def _outline(self, prs, slide, *, idx, total, doc):
        s = prs.slides.add_slide(self._lay(prs, _IDX_BLANK))
        _bg(s)
        T = _T
        # v13: pure-language outline title
        _tb1(s, _S.pick(_S.OUTLINE_TITLE, doc.lang),
             Inches(0), Inches(0.35), Inches(13.333), Inches(0.6),
             Pt(22), T.TEXT, T.LAT_SERIF, T.EA_SERIF, bold=True, align=PP_ALIGN.CENTER)
        _line(s, Inches(0.8), Inches(1.0), Inches(12.5), Inches(1.0), T.RULE, Pt(1))

        n = len(slide.bullets)
        if n <= 8:
            row_h = Inches(0.50)
            font_pt = 16
        elif n <= 12:
            row_h = Inches(0.42)
            font_pt = 14
        else:
            row_h = Inches(0.35)
            font_pt = 12

        content_start_y = Inches(1.1)
        for i, bul in enumerate(slide.bullets):
            y = content_start_y + i * row_h
            # Number: bold dim serif, left-aligned at 2.0" from left
            _tb1(s, f"{i+1:02d}", Inches(2.0), y, Inches(0.6), row_h,
                 Pt(font_pt), T.TEXT_DIM, T.LAT_SERIF, T.EA_SERIF, bold=True)
            # Chapter name: sans body, ~8" wide
            _tb1(s, bul, Inches(2.65), y, Inches(8.7), row_h,
                 Pt(font_pt), T.TEXT, T.LAT_SANS, T.EA_SANS, wrap=True)
            if i < n - 1:
                sep_y = y + row_h - Inches(0.04)
                _line(s, Inches(2.0), sep_y, Inches(11.3), sep_y, T.RULE, Pt(0.3))

        _footer(s, idx, total, chapter="", lang=doc.lang)
        _notes(s, slide.notes)

    def _divider(self, prs, slide, *, idx, total, doc):
        s = prs.slides.add_slide(self._lay(prs, _IDX_BLANK))
        _bg(s)
        T = _T
        num_label = f"§{self._ch + 1}"

        # Section number — centered, gray, 18pt
        _tb1(s, num_label,
             Inches(0), Inches(2.2), Inches(13.333), Inches(0.55),
             Pt(18), T.TEXT_DIM, T.LAT_SERIF, T.EA_SERIF, align=PP_ALIGN.CENTER)

        # Chapter title — centered, large serif bold, 32pt
        _tb1(s, slide.title,
             Inches(1.5), Inches(2.85), Inches(10.333), Inches(1.3),
             Pt(32), T.TEXT, T.LAT_SERIF, T.EA_SERIF, bold=True,
             align=PP_ALIGN.CENTER, wrap=True)

        # Thin rule — centered, 4" wide
        rule_w = Inches(4.0)
        rule_x = (_T.W - rule_w) / 2
        _line(s, rule_x, Inches(4.25), rule_x + rule_w, Inches(4.25), T.RULE, Pt(1))

        # Localized subtitle: "第N章 · {heading}" or "Chapter N · {heading}"
        ch_n = self._ch + 1
        if doc.lang == "zh":
            zh_n = _ZH_NUMS[ch_n - 1] if ch_n <= len(_ZH_NUMS) else str(ch_n)
            subtitle = f"第 {zh_n} 章  ·  {slide.title}"
        else:
            subtitle = f"Chapter {ch_n}  ·  {slide.title}"

        _tb1(s, subtitle,
             Inches(1.5), Inches(4.4), Inches(10.333), Inches(0.5),
             Pt(14), T.TEXT_DIM, T.LAT_SANS, T.EA_SANS, italic=True,
             align=PP_ALIGN.CENTER)

        _footer(s, idx, total, chapter=slide.title, lang=doc.lang)
        _notes(s, slide.notes)

    def _outline_grouped(self, prs, slide, *, idx, total, doc):
        """Outline slide with 4-5 high-level sections (v7). v13: pure-lang title, fixed number column."""
        s = prs.slides.add_slide(self._lay(prs, _IDX_BLANK))
        _bg(s)
        T = _T
        # v13: pick pure-language outline title
        outline_title = _S.pick(_S.OUTLINE_TITLE, doc.lang)
        _tb1(s, outline_title,
             Inches(0), Inches(0.25), Inches(13.333), Inches(0.6),
             Pt(22), T.TEXT, T.LAT_SERIF, T.EA_SERIF, bold=True, align=PP_ALIGN.CENTER)
        _line(s, Inches(0.8), Inches(0.9), Inches(12.5), Inches(0.9), T.RULE, Pt(1))

        n = len(slide.bullets)
        row_h = Inches(0.9)
        content_start_y = Inches(1.05)

        # Parse takeaways from caption field (newline-separated)
        takeaways = slide.caption.split("\n") if slide.caption else []

        for i, name in enumerate(slide.bullets):
            y = content_start_y + i * row_h
            takeaway = takeaways[i] if i < len(takeaways) else ""

            # v13: number column wider (0.8") so "01" fits on one line; use single digit format
            num_str = f"{i+1:02d}"
            _tb1(s, num_str, Inches(2.0), y, Inches(0.8), Inches(0.55),
                 Pt(18), T.TEXT_DIM, T.LAT_SERIF, T.EA_SERIF, bold=True)
            # Section name: bold sans, shifted right to accommodate wider number column
            _tb1(s, name, Inches(2.9), y, Inches(8.4), Inches(0.48),
                 Pt(20), T.TEXT, T.LAT_SANS, T.EA_SANS, bold=True, wrap=True)
            # Takeaway: italic gray, smaller
            if takeaway:
                _tb1(s, takeaway, Inches(2.9), y + Inches(0.48), Inches(8.4), Inches(0.36),
                     Pt(13), T.TEXT_DIM, T.LAT_SANS, T.EA_SANS, italic=True, wrap=True)

            # Separator (not after last item)
            if i < n - 1:
                sep_y = y + row_h - Inches(0.04)
                _line(s, Inches(2.0), sep_y, Inches(11.3), sep_y, T.RULE, Pt(0.3))

        _footer(s, idx, total, chapter="", lang=doc.lang)
        _notes(s, slide.notes)

    def _section_divider(self, prs, slide, *, idx, total, doc):
        """Section divider slide v8: left anchor block + right key-points card.

        Layout (16:9, 13.33"×7.5"):
          Left block  (x=0.5"–4.3"): §N, section title, thin rule, takeaway
          Vertical separator at x=4.5" (y=1.0"–6.5")
          Right block (x=5.0"–12.5"): eyebrow label, rounded-rect card, bullets
          Nav hint bottom-right, footer as usual.

        v12: uses slide.group_idx for the §N label (hierarchical §<group>).
        """
        s = prs.slides.add_slide(self._lay(prs, _IDX_BLANK))
        _bg(s)
        T = _T
        # v12: prefer slide.group_idx (set by SlidePlanner); fall back to counter
        sec_num = slide.group_idx if slide.group_idx > 0 else (self._sec_idx + 1)

        # ── LEFT BLOCK ──────────────────────────────────────────────────────────

        _tb1(s, f"§{sec_num}",
             Inches(0.5), Inches(0.8), Inches(3.8), Inches(0.55),
             Pt(22), _GRAY88, T.LAT_SERIF, T.EA_SERIF,
             bold=False, align=PP_ALIGN.LEFT)

        _tb1(s, slide.title,
             Inches(0.5), Inches(1.6), Inches(3.8), Inches(2.2),
             Pt(28), T.TEXT, T.LAT_SERIF, T.EA_SERIF, bold=True,
             align=PP_ALIGN.LEFT, wrap=True)

        _line(s, Inches(0.5), Inches(4.0), Inches(2.0), Inches(4.0),
              T.RULE, Pt(1))

        if slide.caption:
            _tb1(s, slide.caption,
                 Inches(0.5), Inches(4.3), Inches(3.8), Inches(1.8),
                 Pt(12), T.TEXT_DIM, T.LAT_SANS, T.EA_SANS,
                 italic=True, align=PP_ALIGN.LEFT, wrap=True)

        _line(s, Inches(4.5), Inches(1.0), Inches(4.5), Inches(6.5),
              T.RULE, Pt(1))

        # ── RIGHT BLOCK ─────────────────────────────────────────────────────────

        # v13: pure-language eyebrow label
        eyebrow = _S.pick(_S.SECTION_EYEBROW, doc.lang)
        _tb1(s, eyebrow,
             Inches(5.0), Inches(1.0), Inches(7.5), Inches(0.38),
             Pt(11), _GRAY88, T.LAT_SANS, T.EA_SANS,
             align=PP_ALIGN.LEFT)

        card = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.0), Inches(1.5), Inches(7.5), Inches(4.5)
        )
        card.adjustments[0] = 0.05
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        card.line.color.rgb = RGBColor(0xE5, 0xE5, 0xE5)
        card.line.width = Pt(1)
        card.text_frame.text = ""

        bullets = slide.bullets[:7] if slide.bullets else []   # v13: was 5
        n_bullets = len(bullets)

        if n_bullets > 0:
            card_top    = 1.5
            card_bottom = 6.0
            card_h      = card_bottom - card_top
            padding     = 0.4
            usable_h    = card_h - 2 * padding
            row_h       = usable_h / max(n_bullets, 1)

            for i, bul in enumerate(bullets):
                by = Inches(card_top + padding + i * row_h + row_h * 0.15)
                marker = _MARKERS[i] if i < len(_MARKERS) else "▸"
                _tb1(s, marker,
                     Inches(5.4), by, Inches(0.45), Inches(row_h * 0.7),
                     Pt(14), _GRAY88, T.LAT_SANS, T.EA_SANS,
                     align=PP_ALIGN.LEFT)
                _tb1(s, bul,
                     Inches(5.9), by, Inches(6.35), Inches(row_h * 0.7),
                     Pt(16), T.TEXT, T.LAT_SANS, T.EA_SANS,
                     align=PP_ALIGN.LEFT, wrap=True)

        nav_text = _S.pick(_S.NAV_HINT, doc.lang)
        _tb1(s, nav_text,
             Inches(11.0), Inches(6.3), Inches(1.8), Inches(0.35),
             Pt(9), _GRAY88, T.LAT_SANS, T.EA_SANS,
             italic=True, align=PP_ALIGN.RIGHT)

        # Chapter label for footer (use parent section name)
        _footer(s, idx, total, chapter=slide.title, lang=doc.lang)
        _notes(s, slide.notes)

    def _closing_rich(self, prs, slide, *, idx, total, doc):
        """Rich closing slide with 5-7 bullets + takeaway (v7). v13: pure-language title."""
        s = prs.slides.add_slide(self._lay(prs, _IDX_BLANK))
        _bg(s)
        T = _T

        # v13: pure-language closing title
        _tb1(s, _S.pick(_S.CLOSING_TITLE, doc.lang),
             Inches(0.7), Inches(0.15), Inches(12.0), Inches(0.45),
             Pt(12), T.TEXT_DIM, T.LAT_SERIF, T.EA_SERIF, wrap=True)
        _line(s, Inches(0.7), Inches(0.63), Inches(12.6), Inches(0.63), T.RULE, Pt(0.75))

        body_top = Inches(0.9)
        row_h = Inches(0.62)
        n_bullets = min(len(slide.bullets), 7)
        for i, bul in enumerate(slide.bullets[:n_bullets]):
            by = body_top + i * row_h
            marker = _MARKERS[i] if i < len(_MARKERS) else "▸"
            _tb1(s, marker, Inches(0.7), by, Inches(0.45), row_h,
                 Pt(14), T.TEXT_DIM, T.LAT_SANS, T.EA_SANS)
            _tb1(s, bul, Inches(1.2), by, Inches(11.6), row_h,
                 Pt(16), T.TEXT, T.LAT_SANS, T.EA_SANS, wrap=True)

        # Separator above takeaway
        sep_y = body_top + n_bullets * row_h + Inches(0.1)
        _line(s, Inches(0.7), sep_y, Inches(12.6), sep_y, T.RULE, Pt(0.5))

        # Takeaway
        if slide.caption:
            _tb1(s, f"→ {slide.caption}",
                 Inches(0.7), sep_y + Inches(0.12), Inches(11.6), Inches(0.6),
                 Pt(14), T.TEXT_DIM, T.LAT_SANS, T.EA_SANS, bold=True, italic=True, wrap=True)

        _footer(s, idx, total, chapter="", lang=doc.lang)
        _notes(s, slide.notes)

    def _bullets(self, prs, slide, *, idx, total, doc):
        s = prs.slides.add_slide(self._lay(prs, _IDX_BLANK))
        _bg(s)
        T = _T
        # Use heading→index map so §N is correct even when no divider fired
        chi = self._ch_idx.get(slide.title, max(self._ch, 0))
        is_closing = slide.kind == "closing"

        if is_closing:
            # v13: pure-language closing label
            sec_label = _S.pick(_S.CLOSING_TITLE, doc.lang)
        else:
            kw = "要点" if doc.lang == "zh" else "Key Insights"
            # v12: use hierarchical §group.chapter label when available
            if slide.group_idx > 0 and slide.chapter_in_group > 0:
                num_label = f"§{slide.group_idx}.{slide.chapter_in_group}"
            else:
                num_label = f"§{chi + 1}"
            sec_label = f"{num_label}  ·  {slide.title}  ·  {kw}"

        _tb1(s, sec_label,
             Inches(0.7), Inches(0.15), Inches(12.0), Inches(0.45),
             Pt(12), T.TEXT_DIM, T.LAT_SERIF, T.EA_SERIF, wrap=True)
        _line(s, Inches(0.7), Inches(0.63), Inches(12.6), Inches(0.63), T.RULE, Pt(0.75))

        if slide.bullets:
            body_top = Inches(0.9)
            row_h = Inches(0.62)
            for i, bul in enumerate(slide.bullets):
                by = body_top + i * row_h
                marker = _MARKERS[i] if i < len(_MARKERS) else "▸"
                _tb1(s, marker, Inches(0.7), by, Inches(0.45), row_h,
                     Pt(14), T.TEXT_DIM, T.LAT_SANS, T.EA_SANS)
                _tb1(s, bul, Inches(1.2), by, Inches(11.6), row_h,
                     Pt(16), T.TEXT, T.LAT_SANS, T.EA_SANS, wrap=True)

        _footer(s, idx, total, chapter=self._cur_chapter, lang=doc.lang)
        _notes(s, slide.notes)

    def _figure(self, prs, slide, *, idx, total, doc):
        s = prs.slides.add_slide(self._lay(prs, _IDX_BLANK))
        _bg(s)
        T = _T
        chi = self._ch_idx.get(self._cur_chapter, max(self._ch, 0))

        # Parse figure number from slide.title (e.g. "图 1: CAFE 相调控...")
        fig_num = _parse_fig_num(slide.title)
        fig_label = f"Fig. {fig_num}" if fig_num else "Fig."
        # v12: use hierarchical §group.chapter label when available
        if slide.group_idx > 0 and slide.chapter_in_group > 0:
            sec_label = f"§{slide.group_idx}.{slide.chapter_in_group}"
        else:
            sec_label = f"§{chi + 1}"
        header_text = f"{sec_label}  ·  {fig_label}  ·  {_short_title(slide.title, 50)}"
        _tb1(s, header_text,
             Inches(0.7), Inches(0.15), Inches(12.0), Inches(0.45),
             Pt(12), T.TEXT_DIM, T.LAT_SERIF, T.EA_SERIF, wrap=True)
        _line(s, Inches(0.7), Inches(0.63), Inches(12.6), Inches(0.63), T.RULE, Pt(0.75))

        # Side-by-side layout: alternate image-left / image-right
        img_on_left = (self._fig_idx % 2 == 0)

        slide_w = 13.333
        margin = 0.5
        mid = slide_w / 2.0
        half_w = mid - margin          # ≈6.166"

        body_top = 0.9
        body_h   = 7.5 - body_top - 0.45

        img_area_x  = margin if img_on_left else mid
        text_area_x = mid    if img_on_left else margin

        img_max_w = Inches(half_w)
        img_max_h = Inches(body_h - 0.1)
        img_top   = Inches(body_top + 0.1)

        for ip in slide.image_paths:
            if not ip.exists(): continue
            pic = s.shapes.add_picture(str(ip), Inches(img_area_x), img_top, img_max_w)
            if pic.height > img_max_h:
                r = img_max_h / pic.height
                pic.height = img_max_h
                pic.width  = int(pic.width * r)
            if pic.width > img_max_w:
                r = img_max_w / pic.width
                pic.width  = img_max_w
                pic.height = int(pic.height * r)
            half_center_x = Inches(img_area_x) + (img_max_w - pic.width) // 2
            half_center_y = img_top + (img_max_h - pic.height) // 2
            pic.left = int(half_center_x)
            pic.top  = int(half_center_y)
            break

        tx = Inches(text_area_x + 0.2)
        tw = Inches(half_w - 0.3)

        _tb1(s, slide.title,
             tx, Inches(body_top + 0.1), tw, Inches(0.7),
             Pt(14), T.TEXT, T.LAT_SERIF, T.EA_SERIF, bold=True, wrap=True)
        _line(s, tx, Inches(body_top + 0.85), tx + tw, Inches(body_top + 0.85), T.RULE, Pt(0.5))

        _tb1(s, _S.pick(_S.FIG_EYEBROW, doc.lang),
             tx, Inches(body_top + 1.0), tw, Inches(0.34),
             Pt(13), T.TEXT_DIM, T.LAT_SANS, T.EA_SANS, bold=True)
        _line(s, tx, Inches(body_top + 1.40), tx + tw, Inches(body_top + 1.40),
              T.RULE, Pt(0.3))

        observations = slide.observations or ((slide.deep_observation,) if slide.deep_observation else ())
        obs_row_h = Inches(0.70)
        obs_area_top = Inches(body_top + 1.52)
        obs_area_h = Inches(body_h - 1.62)
        for j, obs_pt in enumerate(observations[:3]):
            oy = obs_area_top + j * obs_row_h
            if oy + obs_row_h > obs_area_top + obs_area_h:
                break
            _tb1(s, "◇", tx, oy, Inches(0.30), obs_row_h,
                 Pt(13), T.TEXT_FAINT, T.LAT_SANS, T.EA_SANS)
            _tb1(s, obs_pt, tx + Inches(0.30), oy, tw - Inches(0.30), obs_row_h,
                 Pt(13), T.TEXT_DIM, T.LAT_SANS, T.EA_SANS, italic=True, wrap=True)

        _footer(s, idx, total, chapter=self._cur_chapter, lang=doc.lang)
        _notes(s, slide.notes)

    def _combined(self, prs, slide, *, idx, total, doc):
        """Combined slide: image on one side, bullets + deep observation on other.

        Alternates image-left / image-right by fig_idx for visual variety.
        Layout (image-left example):
          ┌────────────┐  ┌──────────────────────────────────────┐
          │            │  │  要 点 (small-caps 12pt gray)          │
          │  [IMAGE]   │  │  ❶ bullet 1                          │
          │            │  │  ❷ bullet 2                          │
          │            │  │  ❸ bullet 3                          │
          │            │  │                                      │
          │            │  │  深 度 观 察 (small-caps 12pt gray)   │
          │            │  │  italic one-liner                    │
          └────────────┘  └──────────────────────────────────────┘
        """
        s = prs.slides.add_slide(self._lay(prs, _IDX_BLANK))
        _bg(s)
        T = _T
        chi = self._ch_idx.get(slide.title, max(self._ch, 0))

        fig_num = _parse_fig_num(slide.caption or "") or str(self._fig_idx + 1)
        fig_label = f"Fig. {fig_num}"
        caption_short = _short_title(slide.caption or slide.title, 55)
        # v12: use hierarchical §group.chapter label when available
        if slide.group_idx > 0 and slide.chapter_in_group > 0:
            sec_label = f"§{slide.group_idx}.{slide.chapter_in_group}"
        else:
            sec_label = f"§{chi + 1}"
        # v13: reduced header breadcrumbs (sec_label + fig_label + caption) — removed chapter title
        header_text = f"{sec_label}  ·  {fig_label}  ·  {caption_short}"
        _tb1(s, header_text,
             Inches(0.7), Inches(0.15), Inches(12.0), Inches(0.45),
             Pt(12), T.TEXT_DIM, T.LAT_SERIF, T.EA_SERIF, wrap=True)
        _line(s, Inches(0.7), Inches(0.63), Inches(12.6), Inches(0.63), T.RULE, Pt(0.75))

        # ── layout geometry ───────────────────────────────────────────────────
        img_on_left = (self._fig_idx % 2 == 0)

        slide_w  = 13.333
        margin   = 0.4
        gap      = 0.3        # gap between image and text panes
        body_top = 0.82       # top of body area (below rule + small clearance)
        body_h   = 7.5 - body_top - 0.45  # ≈ 6.23"

        # Image pane: 5.5" wide; text pane: remainder
        img_pane_w  = 5.8
        text_pane_w = slide_w - 2 * margin - img_pane_w - gap   # ≈ 6.033"

        if img_on_left:
            img_x  = margin
            text_x = margin + img_pane_w + gap
        else:
            text_x = margin
            img_x  = margin + text_pane_w + gap

        # ── image pane ────────────────────────────────────────────────────────
        img_max_w = Inches(img_pane_w)
        img_max_h = Inches(body_h - 0.1)
        img_top   = Inches(body_top + 0.05)

        for ip in slide.image_paths:
            if not ip.exists():
                continue
            pic = s.shapes.add_picture(str(ip), Inches(img_x), img_top, img_max_w)
            if pic.height > img_max_h:
                r = img_max_h / pic.height
                pic.height = img_max_h
                pic.width  = int(pic.width * r)
            if pic.width > img_max_w:
                r = img_max_w / pic.width
                pic.width  = img_max_w
                pic.height = int(pic.height * r)
            # Center within pane
            pic.left = int(Inches(img_x) + (img_max_w - pic.width) // 2)
            pic.top  = int(img_top + (img_max_h - pic.height) // 2)
            break

        # ── text pane ─────────────────────────────────────────────────────────
        tx = Inches(text_x)
        tw = Inches(text_pane_w - 0.1)
        ty = Inches(body_top + 0.1)
        slide_bottom = Inches(body_top + body_h)

        observations = slide.observations or ((slide.deep_observation,) if slide.deep_observation else ())
        has_bullets = bool(slide.bullets)
        has_obs     = bool(observations)

        if has_bullets:
            # ── bullets section (upper part of right pane) ─────────────────
            _tb1(s, _S.pick(_S.COMBINED_KW, doc.lang), tx, ty, tw, Inches(0.35),
                 Pt(12), T.TEXT_DIM, T.LAT_SANS, T.EA_SANS, bold=True)

            bullet_row_h = Inches(0.55)
            bullets_top = ty + Inches(0.38)
            for i, bul in enumerate(slide.bullets):
                by = bullets_top + i * bullet_row_h
                marker = _MARKERS[i] if i < len(_MARKERS) else "▸"
                _tb1(s, marker, tx, by, Inches(0.35), bullet_row_h,
                     Pt(13), T.TEXT_DIM, T.LAT_SANS, T.EA_SANS)
                _tb1(s, bul, tx + Inches(0.38), by, tw - Inches(0.38), bullet_row_h,
                     Pt(14), T.TEXT, T.LAT_SANS, T.EA_SANS, wrap=True)

            n_bullets = len(slide.bullets)
            obs_top = bullets_top + n_bullets * bullet_row_h + Inches(0.15)
        else:
            # ── no bullets: observations fill the full right pane ──────────
            obs_top = ty

        if has_obs:
            # Clamp obs_top so observations leave room for at least one row.
            obs_top = min(obs_top, slide_bottom - Inches(1.9))

            _tb1(s, _S.pick(_S.FIG_EYEBROW, doc.lang), tx, obs_top, tw, Inches(0.32),
                 Pt(13), T.TEXT_DIM, T.LAT_SANS, T.EA_SANS, bold=True)
            _line(s, tx, obs_top + Inches(0.36), tx + tw, obs_top + Inches(0.36),
                  T.RULE, Pt(0.3))

            obs_item_top = obs_top + Inches(0.44)
            obs_area_h   = slide_bottom - obs_item_top - Inches(0.05)
            n_obs = min(len(observations), 3)
            obs_row_h = min(Inches(0.70), int(obs_area_h / max(n_obs, 1)))
            for j, obs_pt in enumerate(observations[:n_obs]):
                oy = obs_item_top + j * obs_row_h
                if oy + obs_row_h > slide_bottom - Inches(0.05):
                    break
                _tb1(s, "◇", tx, oy, Inches(0.30), obs_row_h,
                     Pt(13), T.TEXT_FAINT, T.LAT_SANS, T.EA_SANS)
                _tb1(s, obs_pt, tx + Inches(0.30), oy, tw - Inches(0.30), obs_row_h,
                     Pt(13), T.TEXT_DIM, T.LAT_SANS, T.EA_SANS, italic=True, wrap=True)

        _footer(s, idx, total, chapter=self._cur_chapter, lang=doc.lang)
        _notes(s, slide.notes)

    @staticmethod
    def _lay(prs, idx):
        lays = prs.slide_layouts
        return lays[idx] if idx < len(lays) else lays[0]


# ── design primitives ──────────────────────────────────────────────────────────

def _bg(s) -> None:
    """Cream full-slide rectangle pushed to back of z-order."""
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), _T.W, _T.H)
    bg.fill.solid(); bg.fill.fore_color.rgb = _T.BG
    bg.line.fill.background()
    tree = s.shapes._spTree; el = bg._element
    tree.remove(el); tree.insert(2, el)


def _line(s, x1, y1, x2, y2, color, width):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.line.color.rgb = color; ln.line.width = width


def _footer(s, slide_idx, total, *, chapter: str = "", lang: str = "zh") -> None:
    T = _T
    if chapter:
        footer_prefix = _S.CHAPTER_FOOTER[0] if lang == "zh" else _S.CHAPTER_FOOTER[1]
        _tb1(s, footer_prefix + chapter,
             Inches(0.7), Inches(7.1), Inches(8.0), Inches(0.3),
             Pt(9), T.TEXT_FAINT, T.LAT_SANS, T.EA_SANS, italic=True)
    _tb1(s, f"{slide_idx} of {total}",
         Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
         Pt(9), T.TEXT_FAINT, T.LAT_SANS, T.EA_SANS)


def _notes(s, text):
    if text:
        s.notes_slide.notes_text_frame.text = text


def _tb1(s, text, left, top, w, h, fsize, fcolor, lat_font, ea_font,
         bold=False, italic=False, wrap=False, align=None):
    """Add a single-paragraph textbox."""
    tb = s.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    if wrap: tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    _run_style(r, fsize, bold, fcolor, lat_font, ea_font, italic=italic)
    if align is not None: p.alignment = align


def _run_style(r, size, bold, color, lat_font, ea_font, italic=False):
    r.font.size = size; r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = lat_font
    _ea(r, ea_font)


def _ea(run, font_name):
    """Inject <a:ea typeface="..."/> for CJK rendering."""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_name)


def _parse_fig_num(title: str) -> str:
    """Extract figure number from slide title like '图 1: ...' or 'Fig. 2 ...'"""
    m = re.search(r'(?:图|Fig\.?)\s*(\d+)', title, re.IGNORECASE)
    return m.group(1) if m else ""


def _short_title(title: str, maxlen: int) -> str:
    """Truncate title, stripping the 'Fig N: ' prefix first."""
    # Remove leading "图 N: " or "Fig. N: " portion
    clean = re.sub(r'^(?:图|Fig\.?)\s*\d+[:\s]*', '', title, flags=re.IGNORECASE).strip()
    return clean[:maxlen] + ("…" if len(clean) > maxlen else "")


RENDERERS["pptx"] = PptxRenderer
