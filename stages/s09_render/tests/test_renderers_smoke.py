from pathlib import Path

import pytest
from docx import Document as DocxDocument
from PIL import Image

from stages.s09_render.model import Document, Chapter, Paragraph, FigureBlock, TableBlock
from stages.s09_render.renderers import RENDERERS
import stages.s09_render.renderers.docx  # noqa: F401 — triggers RENDERERS["docx"] registration
import stages.s09_render.renderers.html  # noqa: F401 — triggers RENDERERS["html"] registration
import stages.s09_render.renderers.pdf   # noqa: F401 — triggers RENDERERS["pdf"] registration
import stages.s09_render.renderers.pptx  # noqa: F401 — triggers RENDERERS["pptx"] registration


@pytest.fixture
def one_image(tmp_path: Path) -> Path:
    p = tmp_path / "tiny.jpg"
    Image.new("RGB", (100, 50), "red").save(p)
    return p


def _make_doc(one_image: Path) -> Document:
    return Document(
        paper_title="Smoke Test Paper",
        lang="zh",
        chapters=(
            Chapter(heading="引言", level=1, blocks=(
                Paragraph(text="这是引言的第一段。"),
                Paragraph(text="第二段提到 Fig. 1 的内容。"),
                FigureBlock(fig_id="Fig. 1", label="图 1",
                            image_paths=(one_image,),
                            caption="第一张图", deep_observation="观察"),
            )),
        ),
    )


def test_docx_renderer_produces_readable_file(tmp_path: Path, one_image: Path):
    doc = _make_doc(one_image)
    out = tmp_path / "preview.docx"
    RENDERERS["docx"]().render(doc, out)
    assert out.exists() and out.stat().st_size > 4000
    d = DocxDocument(out)
    text = "\n".join(p.text for p in d.paragraphs)
    assert "Smoke Test Paper" in text
    assert "引言" in text
    assert "图 1. 第一张图" in text
    assert len(d.inline_shapes) == 1


def test_html_renderer_self_contained_base64(tmp_path: Path, one_image: Path):
    doc = _make_doc(one_image)
    out = tmp_path / "preview.html"
    RENDERERS["html"]().render(doc, out)
    assert out.exists()
    html = out.read_text(encoding="utf-8")
    assert "Smoke Test Paper" in html
    assert "引言" in html
    # v2 (2026-06): caption label is wrapped in <span class="fig-tag"> so
    # the literal "图 1. 第一张图" is split across tags. Verify both parts.
    assert '<span class="fig-tag">图 1.</span>' in html
    assert "第一张图" in html
    # Base64 embedded image — no external file refs
    assert 'src="data:image/' in html
    assert 'src="/tmp' not in html  # absolute paths must NOT leak


def test_pdf_renderer_produces_valid_pdf_file(tmp_path: Path, one_image: Path):
    doc = _make_doc(one_image)
    out = tmp_path / "preview.pdf"
    RENDERERS["pdf"]().render(doc, out)
    assert out.exists()
    assert out.read_bytes()[:5] == b"%PDF-"
    assert out.stat().st_size > 10_000  # cover page + 1 image


def test_pptx_renderer_produces_valid_deck(tmp_path: Path, one_image: Path):
    from pptx import Presentation
    doc = _make_doc(one_image)
    out = tmp_path / "preview.pptx"
    RENDERERS["pptx"]().render(doc, out)
    assert out.exists() and out.stat().st_size > 10_000

    prs = Presentation(str(out))
    n = len(prs.slides)
    # Minimum: title + outline + at least one content + closing = 4
    assert n >= 4
    # First slide is the title
    title_shape = prs.slides[0].shapes.title
    assert title_shape is not None
    assert "Smoke Test Paper" in title_shape.text


def test_pptx_renderer_accepts_user_template(tmp_path: Path, one_image: Path):
    """User-supplied template.pptx overrides the default base presentation."""
    from pptx import Presentation
    # Generate a "template" by running once without it
    doc = _make_doc(one_image)
    seed = tmp_path / "seed.pptx"
    RENDERERS["pptx"]().render(doc, seed)
    # Now render WITH it as a template
    out = tmp_path / "out.pptx"
    RENDERERS["pptx"](template_path=seed).render(doc, out)
    assert out.exists() and out.stat().st_size > 10_000
    prs = Presentation(str(out))
    # Sanity: still has slides
    assert len(prs.slides) >= 4


def test_pptx_renderer_layout_fallback_when_template_has_few_layouts(
        tmp_path: Path, one_image: Path):
    """If a template has fewer layouts than expected, layout selection falls back."""
    from stages.s09_render.renderers.pptx import PptxRenderer
    from pptx import Presentation
    prs = Presentation()
    layout = PptxRenderer._lay(prs, 999)  # out of range
    assert layout is prs.slide_layouts[0]


def test_table_block_parse_round_trip():
    """TableBlock parsing correctly handles headers, separator, and data rows."""
    from stages.s09_render.builder import DocumentBuilder
    from stages.s09_render.model import TableBlock
    builder = DocumentBuilder(lang="zh", paper_title="T")
    md_table = "| A | B | C |\n|---|---|---|\n| 1 | 2 | 3 |\n| 4 | 5 | 6 |"
    block = builder._parse_md_table(md_table)
    assert isinstance(block, TableBlock)
    assert block.headers == ("A", "B", "C")
    assert len(block.rows) == 2
    assert block.rows[0] == ("1", "2", "3")


def test_docx_table_block_renders(tmp_path: Path, one_image: Path):
    """TableBlock in a chapter renders as a Word table (not pipe-text) in DOCX."""
    doc = Document(
        paper_title="Table Test Paper",
        lang="zh",
        chapters=(
            Chapter(heading="数据对比", level=1, blocks=(
                Paragraph(text="以下是比较数据。"),
                TableBlock(
                    headers=("材料", "Wrec (J/cm³)", "η (%)"),
                    rows=(
                        ("样品A", "8.6", "85"),
                        ("样品B", "6.2", "78"),
                    ),
                ),
            )),
        ),
    )
    out = tmp_path / "preview.docx"
    RENDERERS["docx"]().render(doc, out)
    d = DocxDocument(out)
    # Should have at least 1 Word table
    assert len(d.tables) >= 1
    # Header text should be in the first table
    first_table_text = " ".join(
        cell.text for row in d.tables[0].rows for cell in row.cells
    )
    assert "材料" in first_table_text
    assert "样品A" in first_table_text


def test_html_table_block_renders(tmp_path: Path, one_image: Path):
    """TableBlock renders as <table> in HTML output."""
    doc = Document(
        paper_title="HTML Table Test",
        lang="zh",
        chapters=(
            Chapter(heading="比较", level=1, blocks=(
                TableBlock(
                    headers=("系统", "性能"),
                    rows=(("AFE-A", "高"),),
                ),
            )),
        ),
    )
    out = tmp_path / "preview.html"
    RENDERERS["html"]().render(doc, out)
    html = out.read_text(encoding="utf-8")
    assert "<table" in html
    assert "系统" in html
    assert "AFE-A" in html


def test_builder_split_paragraphs_with_table():
    """_split_paragraphs yields TableBlock for markdown table text."""
    from stages.s09_render.builder import DocumentBuilder
    builder = DocumentBuilder(lang="zh", paper_title="T")
    body = "普通段落。\n\n| H1 | H2 |\n|---|---|\n| v1 | v2 |\n\n另一段落。"
    blocks = list(builder._split_paragraphs(body))
    # Should be: Paragraph, TableBlock, Paragraph
    from stages.s09_render.model import Paragraph, TableBlock
    assert len(blocks) == 3
    assert isinstance(blocks[0], Paragraph)
    assert isinstance(blocks[1], TableBlock)
    assert isinstance(blocks[2], Paragraph)
    assert blocks[1].headers == ("H1", "H2")


def test_pptx_renderer_adds_footer_with_slide_numbers(tmp_path: Path, one_image: Path):
    """Footer text 'N of total' appears on non-title slides."""
    doc = _make_doc(one_image)
    out = tmp_path / "fp.pptx"
    RENDERERS["pptx"]().render(doc, out)
    from pptx import Presentation
    prs = Presentation(str(out))
    # Slide 2 onwards should have a footer textbox with "of total" pattern
    total = len(prs.slides)
    found_footer = False
    all_slides = list(prs.slides)
    for slide in all_slides[1:]:  # skip title slide (index 0)
        for shape in slide.shapes:
            if shape.has_text_frame and f"of {total}" in shape.text_frame.text:
                found_footer = True
                break
        if found_footer:
            break
    assert found_footer, "expected footer 'N of total' on a non-title slide"
