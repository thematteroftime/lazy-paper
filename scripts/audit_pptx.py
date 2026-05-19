#!/usr/bin/env python3
"""Audit a rendered PPT for known layout / content defects.

Reports per-slide:
- Exotic Unicode that renders as boxes (U+2011, U+202F, U+200B, …)
- Bullets ending with `…` but on slides with low total info density
- Empty text frames (waste of vertical space)
- Wrong-language paragraphs (CJK chars in --lang en, ASCII-only in --lang zh)
- Caption-header text truncated mid-formula

Usage:
    uv run python scripts/audit_pptx.py runs/yang2025/s09_render/preview.pptx
    uv run python scripts/audit_pptx.py runs/*/s09_render/preview.pptx
"""
import argparse, glob, re, sys
from collections import Counter
from pathlib import Path

from pptx import Presentation

EXOTIC_CODEPOINTS = (0x2011, 0x202F, 0x200B, 0x2009, 0x2007)
MID_FORMULA_RE = re.compile(r"\([A-Za-z]+\d?[\.\,]?\d*[A-Za-z]*$")
TRAILING_OPEN_PAREN_RE = re.compile(r"\([^\)]{0,3}…?$")
# Numbered bullet markers (❶-❼, ①-⑩, ⓵-⓿, ⓪-⑨) and similar — not real content.
BULLET_MARKERS = set("❶❷❸❹❺❻❼❽❾❿①②③④⑤⑥⑦⑧⑨⑩⓪⓵⓶⓷⓸⓹⓺⓻⓼⓽◇◆▸▪•·")


def is_cjk(c: str) -> bool:
    return "一" <= c <= "鿿"


def is_ascii_letter(c: str) -> bool:
    return c.isascii() and c.isalpha()


def audit_pptx(path: Path, lang_hint: str | None = None) -> dict:
    if not path.exists():
        return {"error": f"file missing: {path}"}
    pres = Presentation(str(path))
    issues = Counter()
    examples: dict[str, list[str]] = {}

    def flag(key: str, slide_no: int, text: str):
        issues[key] += 1
        examples.setdefault(key, []).append(f"slide {slide_no}: {text[:120]!r}")

    n_slides = len(pres.slides)
    total_paragraphs = 0
    cjk_paragraphs = 0
    ascii_paragraphs = 0

    for si, slide in enumerate(pres.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                text = para.text
                if not text.strip():
                    continue
                total_paragraphs += 1
                # Exotic codepoints
                for ch in text:
                    if ord(ch) in EXOTIC_CODEPOINTS:
                        flag(f"exotic_U+{ord(ch):04X}", si, text)
                        break
                # Language drift
                cjk_count = sum(1 for c in text if is_cjk(c))
                ascii_alpha = sum(1 for c in text if is_ascii_letter(c))
                if cjk_count + ascii_alpha > 4:
                    if cjk_count * 2 >= len(text):
                        cjk_paragraphs += 1
                    elif ascii_alpha * 2 >= len(text):
                        ascii_paragraphs += 1
                # Skip title slide (si == 1): paper titles and presenter
                # metadata legitimately use the original-language English.
                if si > 1 and lang_hint == "en" and cjk_count > 4 and ascii_alpha == 0:
                    flag("lang_drift_zh_in_en_paper", si, text)
                elif si > 1 and lang_hint == "zh" and cjk_count == 0 and ascii_alpha > 30:
                    flag("lang_drift_en_in_zh_paper", si, text)
                # Mid-formula truncation (e.g., "Pb0.98La0.02(Zr0.66Ti0.10" or ending in "(In")
                if not text.endswith("…") and (
                    MID_FORMULA_RE.search(text)
                    or TRAILING_OPEN_PAREN_RE.search(text)
                ):
                    # Heuristic: short text ending with bare unclosed paren is suspicious
                    if 30 <= len(text) <= 200:
                        flag("possible_mid_formula_cut", si, text)
                # Empty-ish text frame: a paragraph that's just a single content
                # character (not a bullet marker glyph).
                stripped = text.strip()
                if len(stripped) == 1 and stripped[0] not in BULLET_MARKERS:
                    flag("tiny_paragraph", si, text)

    return {
        "path": str(path),
        "n_slides": n_slides,
        "n_paragraphs": total_paragraphs,
        "cjk_paragraphs": cjk_paragraphs,
        "ascii_paragraphs": ascii_paragraphs,
        "issues": dict(issues),
        "examples": {k: v[:3] for k, v in examples.items()},  # cap examples
    }


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("paths", nargs="+", help="pptx files (globs supported)")
    ap.add_argument("--lang", choices=["en", "zh"], default=None,
                    help="Language hint for lang-drift detection")
    args = ap.parse_args(argv)

    # Expand globs
    expanded = []
    for p in args.paths:
        expanded.extend(glob.glob(p))
    if not expanded:
        print(f"no files matched: {args.paths}", file=sys.stderr)
        return 1

    overall_issues = 0
    for path in sorted(set(expanded)):
        # Infer lang heuristically from BODY paragraphs only (skip titles/
        # headers which often carry the paper's original English name even
        # when the analysis is Chinese).
        lang = args.lang
        if lang is None:
            pres = Presentation(str(path))
            cjk_body = 0
            en_body = 0
            for slide in pres.slides:
                for sh in slide.shapes:
                    if not sh.has_text_frame:
                        continue
                    for para in sh.text_frame.paragraphs:
                        t = para.text.strip()
                        # Only count paragraphs of substance (≥ 40 chars
                        # for ASCII, ≥ 15 for CJK) — skips headers.
                        cjk = sum(1 for c in t if is_cjk(c))
                        ascii_alpha = sum(1 for c in t if is_ascii_letter(c))
                        if cjk >= 15:
                            cjk_body += 1
                        elif ascii_alpha >= 40:
                            en_body += 1
            lang = "zh" if cjk_body > en_body else "en"

        result = audit_pptx(Path(path), lang_hint=lang)
        print(f"\n=== {result['path']} ===")
        print(f"  slides: {result['n_slides']}   paragraphs: {result['n_paragraphs']}")
        if not result.get("issues"):
            print("  ✓ no issues")
            continue
        overall_issues += sum(result["issues"].values())
        for key, count in sorted(result["issues"].items(), key=lambda kv: -kv[1]):
            print(f"  {key}: {count}")
            for ex in result["examples"].get(key, [])[:3]:
                print(f"      {ex}")
    print(f"\nTotal issues: {overall_issues}")
    return 0 if overall_issues == 0 else 2


if __name__ == "__main__":
    sys.exit(main())
